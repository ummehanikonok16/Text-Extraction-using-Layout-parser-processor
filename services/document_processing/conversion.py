import os
import subprocess
from pathlib import Path
import tempfile
import shutil
from typing import Optional
from config.config import Config

class DocumentConverter:
    """Factory class for document conversion operations"""
    
    def __init__(self) -> None:
        self.config = Config.get_instance()
    
    def convert_to_pdf(self, file_path: str) -> str:
        """
        Convert any file type to PDF format if it's not already a PDF.
        Returns the path to the PDF file.
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # If already PDF, return original path
        if file_extension == '.pdf':
            print("File is already in PDF format.")
            return file_path
        
        # Define output PDF path with proper encoding handling
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        output_dir = os.path.dirname(file_path)
        output_pdf_path = os.path.join(output_dir, f"{base_name}_converted.pdf")
        
        print(f"Converting {file_extension} file to PDF...")
        
        try:
            # Office documents (including Excel)
            if file_extension in ['.doc', '.docx', '.ppt', '.pptx', '.xls', '.xlsx', '.odt', '.ods', '.odp']:
                self._convert_office_to_pdf(file_path, output_pdf_path)
            
            # Text-based files
            elif file_extension in ['.txt', '.rtf', '.csv', '.html', '.htm', '.xml', '.json', '.yaml', '.yml', '.py', '.js', '.css', '.java', '.cpp', '.c', '.sql']:
                self._convert_text_to_pdf(file_path, output_pdf_path)
            
            # Image files
            elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif', '.webp', '.svg']:
                self._convert_image_to_pdf(file_path, output_pdf_path)
            
            # Archive files - extract first then convert
            elif file_extension in ['.zip', '.rar', '.7z']:
                self._convert_archive_to_pdf(file_path, output_pdf_path)
            
            # Unknown or unsupported formats
            else:
                print(f"Attempting to convert unsupported format {file_extension}...")
                # Try LibreOffice first (it supports many formats)
                try:
                    self._convert_office_to_pdf(file_path, output_pdf_path)
                except:
                    # If that fails, try to read as text
                    try:
                        self._convert_text_to_pdf(file_path, output_pdf_path)
                    except:
                        # Last resort: create a PDF with file info
                        self._create_fallback_pdf(file_path, output_pdf_path)
            
            if os.path.exists(output_pdf_path):
                print(f"Successfully converted to PDF: {os.path.basename(output_pdf_path)}")
                return output_pdf_path
            else:
                print("Conversion failed, using original file.")
                return file_path
                
        except Exception as e:
            print(f"Error during conversion: {e}")
            print("Using original file for processing.")
            return file_path
    
    def _convert_office_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert Office documents to PDF using multiple methods"""
        # Method 1: Try LibreOffice
        try:
            self._convert_with_libreoffice(input_path, output_path)
            return
        except Exception as e:
            print(f"LibreOffice method failed: {e}")
        
        # Method 2: Try Python libraries for specific formats
        try:
            file_extension = os.path.splitext(input_path)[1].lower()
            
            if file_extension in ['.xlsx', '.xls']:
                self._convert_excel_to_pdf(input_path, output_path)
            elif file_extension in ['.docx', '.doc']:
                self._convert_word_to_pdf(input_path, output_path)
            elif file_extension in ['.pptx', '.ppt']:
                self._convert_powerpoint_to_pdf(input_path, output_path)
            else:
                raise Exception(f"No alternative method for {file_extension}")
                
        except Exception as e:
            print(f"Python library method failed: {e}")
            raise Exception("All conversion methods failed")
    
    def _convert_with_libreoffice(self, input_path: str, output_path: str) -> None:
        """Convert using LibreOffice command line"""
        try:
            # Create a temporary directory to avoid filename issues
            with tempfile.TemporaryDirectory() as temp_dir:
                # Copy file to temp directory with simple name
                temp_input = os.path.join(temp_dir, "input" + os.path.splitext(input_path)[1])
                shutil.copy2(input_path, temp_input)
                
                # Run LibreOffice conversion
                cmd = [
                    'libreoffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', temp_dir, temp_input
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    # Find the generated PDF
                    temp_pdf = os.path.join(temp_dir, "input.pdf")
                    if os.path.exists(temp_pdf):
                        shutil.copy2(temp_pdf, output_path)
                        return
                    else:
                        raise Exception("PDF not generated by LibreOffice")
                else:
                    raise Exception(f"LibreOffice failed: {result.stderr}")
                    
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out")
        except FileNotFoundError:
            raise Exception("LibreOffice not found. Please install LibreOffice.")
    
    def _convert_excel_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert Excel files to PDF using openpyxl and reportlab"""
        try:
            import openpyxl
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            
            # Load the Excel file
            wb = openpyxl.load_workbook(input_path, data_only=True)
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()
            
            # Process each worksheet
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Add sheet title
                elements.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading1']))
                elements.append(Spacer(1, 12))
                
                # Get data from sheet
                data = []
                for row in ws.iter_rows(values_only=True):
                    # Convert None values to empty strings and handle encoding
                    row_data = []
                    for cell in row:
                        if cell is None:
                            row_data.append("")
                        else:
                            # Handle different data types and encoding
                            try:
                                row_data.append(str(cell))
                            except:
                                row_data.append("(unreadable)")
                    data.append(row_data)
                
                # Skip empty sheets
                if not data or all(all(cell == "" for cell in row) for row in data):
                    elements.append(Paragraph("(Empty sheet)", styles['Normal']))
                else:
                    # Create table
                    table = Table(data)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                        ('FONTSIZE', (0, 1), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black)
                    ]))
                    elements.append(table)
                
                elements.append(Spacer(1, 20))
            
            # Build PDF
            doc.build(elements)
            
        except ImportError as e:
            raise Exception(f"Required libraries not installed: {e}. Try: pip install openpyxl reportlab")
        except Exception as e:
            raise Exception(f"Excel to PDF conversion failed: {e}")
    
    def _convert_word_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert Word documents to PDF"""
        try:
            # Try docx2pdf first
            from docx2pdf import convert
            convert(input_path, output_path)
        except ImportError:
            # Try python-docx + reportlab
            try:
                from docx import Document
                from reportlab.platypus import SimpleDocTemplate, Paragraph
                from reportlab.lib.styles import getSampleStyleSheet
                
                doc = Document(input_path)
                
                # Create PDF
                pdf_doc = SimpleDocTemplate(output_path)
                styles = getSampleStyleSheet()
                story = []
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        story.append(Paragraph(para.text, styles['Normal']))
                
                pdf_doc.build(story)
                
            except ImportError:
                raise Exception("Required libraries not installed. Try: pip install python-docx docx2pdf")
    
    def _convert_powerpoint_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert PowerPoint to PDF using python-pptx"""
        try:
            from pptx import Presentation
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            
            prs = Presentation(input_path)
            
            # Create PDF
            doc = SimpleDocTemplate(output_path)
            styles = getSampleStyleSheet()
            story = []
            
            for i, slide in enumerate(prs.slides):
                story.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
                story.append(Spacer(1, 12))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))
                
                story.append(Spacer(1, 20))
            
            doc.build(story)
            
        except ImportError:
            raise Exception("Required libraries not installed. Try: pip install python-pptx")
    
    def _convert_text_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert text file to PDF"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            
            # Read the text file with various encodings
            text = self._read_text_file(input_path)
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Split text into paragraphs
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    # Clean the text for PDF
                    clean_para = para.replace('\n', '<br/>')
                    story.append(Paragraph(clean_para, styles['Normal']))
                    story.append(Spacer(1, 12))
            
            doc.build(story)
            
        except ImportError:
            # Fallback: create simple text-based PDF using basic approach
            self._create_simple_text_pdf(input_path, output_path)
    
    def _read_text_file(self, file_path: str) -> str:
        """Read text file with multiple encoding attempts"""
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii', 'cp1253']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, read as binary and decode with errors='ignore'
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    
    def _convert_archive_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert archive files to PDF by listing contents"""
        try:
            import zipfile
            
            file_extension = os.path.splitext(input_path)[1].lower()
            content_text = f"Archive Contents: {os.path.basename(input_path)}\n\n"
            
            if file_extension == '.zip':
                with zipfile.ZipFile(input_path, 'r') as zip_file:
                    file_list = zip_file.namelist()
                    content_text += "Files in archive:\n"
                    for file_name in file_list:
                        content_text += f"- {file_name}\n"
            
            # Create temporary text file and convert to PDF
            temp_txt = output_path + '.tmp.txt'
            with open(temp_txt, 'w', encoding='utf-8') as f:
                f.write(content_text)
            
            self._convert_text_to_pdf(temp_txt, output_path)
            
            # Clean up temp file
            if os.path.exists(temp_txt):
                os.remove(temp_txt)
                
        except ImportError:
            self._create_fallback_pdf(input_path, output_path)
    
    def _create_fallback_pdf(self, input_path: str, output_path: str) -> None:
        """Create a fallback PDF with basic file information"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            # Get file info
            file_size = os.path.getsize(input_path)
            file_name = os.path.basename(input_path)
            file_ext = os.path.splitext(input_path)[1]
            
            # Write file information
            y = height - 100
            c.drawString(50, y, f"File: {file_name}")
            y -= 30
            c.drawString(50, y, f"Extension: {file_ext}")
            y -= 30
            c.drawString(50, y, f"Size: {file_size} bytes ({file_size/1024:.2f} KB)")
            y -= 30
            c.drawString(50, y, "Note: This file type could not be converted to readable text.")
            y -= 30
            c.drawString(50, y, "Original file will be processed directly by Document AI.")
            
            c.save()
            
        except ImportError:
            # Ultimate fallback - just copy the original file
            print(f"Cannot create PDF for {input_path}. Using original file.")
            raise Exception("PDF creation failed")
    
    def _create_simple_text_pdf(self, input_path: str, output_path: str) -> None:
        """Create a simple PDF from text using basic approach"""
        # This is a fallback method - in practice, you'd want reportlab or similar
        print("Warning: Advanced PDF creation libraries not available. Using original text file.")
        raise Exception("PDF creation libraries not available")
    
    def _convert_image_to_pdf(self, input_path: str, output_path: str) -> None:
        """Convert image to PDF"""
        try:
            from PIL import Image
            
            image = Image.open(input_path)
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            image.save(output_path, 'PDF')
            
        except ImportError:
            # Try alternative method
            try:
                import subprocess
                # Use ImageMagick if available
                cmd = ['convert', input_path, output_path]
                subprocess.run(cmd, check=True, timeout=30)
            except (FileNotFoundError, subprocess.CalledProcessError):
                raise Exception("Image conversion libraries not available")